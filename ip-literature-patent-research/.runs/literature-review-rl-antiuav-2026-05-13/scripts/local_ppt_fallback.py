from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
RUN=Path('ip-literature-patent-research/.runs/literature-review-rl-antiuav-2026-05-13')
DOC=RUN/'强化学习在反无人机系统中的应用-文献综述.docx'
PPT=RUN/'强化学习在反无人机系统中的应用-组会汇报.pptx'
# Extract rough references for final slide
paras=[p.text.strip() for p in Document(DOC).paragraphs if p.text.strip()]
# Theme
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BG=RGBColor(12,18,32); FG=RGBColor(235,242,255); ACC=RGBColor(96,165,250); MUTED=RGBColor(180,190,210)
FONT='Noto Sans CJK SC'
def add_bg(slide):
    rect=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb=BG; rect.line.fill.background(); rect.z_order=0
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(0.12),prs.slide_height)
    bar.fill.solid(); bar.fill.fore_color.rgb=ACC; bar.line.fill.background()
def add_title(slide,title,subtitle=None):
    t=slide.shapes.add_textbox(Inches(0.65),Inches(0.35),Inches(12),Inches(0.65)).text_frame
    t.clear(); p=t.paragraphs[0]; p.text=title; p.font.name=FONT; p.font.size=Pt(28); p.font.bold=True; p.font.color.rgb=FG
    if subtitle:
        s=slide.shapes.add_textbox(Inches(0.68),Inches(1.02),Inches(11.8),Inches(0.35)).text_frame
        s.text=subtitle; s.paragraphs[0].font.name=FONT; s.paragraphs[0].font.size=Pt(13); s.paragraphs[0].font.color.rgb=MUTED
def bullet_slide(title, bullets, section=''):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide,title,section)
    box=slide.shapes.add_textbox(Inches(0.85),Inches(1.55),Inches(11.8),Inches(5.2)); tf=box.text_frame; tf.clear()
    for i,b in enumerate(bullets):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=b; p.level=0
        p.font.name=FONT; p.font.size=Pt(20 if len(b)<26 else 18); p.font.color.rgb=FG; p.space_after=Pt(12)
    return slide
# Cover
slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
tf=slide.shapes.add_textbox(Inches(0.9),Inches(1.8),Inches(11.7),Inches(1.4)).text_frame; tf.clear()
p=tf.paragraphs[0]; p.text='强化学习在反无人机系统中的应用'; p.font.name=FONT; p.font.size=Pt(38); p.font.bold=True; p.font.color.rgb=FG; p.alignment=PP_ALIGN.CENTER
sf=slide.shapes.add_textbox(Inches(1.2),Inches(3.25),Inches(11),Inches(0.8)).text_frame; sf.text='组会汇报 · 基于文献综述整理'; sf.paragraphs[0].font.name=FONT; sf.paragraphs[0].font.size=Pt(22); sf.paragraphs[0].font.color.rgb=ACC; sf.paragraphs[0].alignment=PP_ALIGN.CENTER
sf2=slide.shapes.add_textbox(Inches(1.2),Inches(4.25),Inches(11),Inches(0.5)).text_frame; sf2.text='2026-05-13'; sf2.paragraphs[0].font.name=FONT; sf2.paragraphs[0].font.size=Pt(16); sf2.paragraphs[0].font.color.rgb=MUTED; sf2.paragraphs[0].alignment=PP_ALIGN.CENTER
slides=[
('汇报路线图',['背景与问题定义','强化学习方法谱系','任务场景与系统建模','代表工作与平台指标','挑战、开放问题与结论'],''),
('研究背景与意义',['小型无人机低成本、集群化、智能化趋势明显','传统固定规则/单一传感器方法难适应强对抗','反无人机需要实时感知、决策与协同拦截','RL 适合在复杂动态环境中学习序贯策略'],'第1节 引言与背景'),
('问题定义与关键挑战',['目标：发现、识别、跟踪、干扰或拦截来袭无人机','状态空间包含位置、速度、航迹、威胁等级与资源','动作空间包括机动、火控、干扰功率、协同分配','挑战：不完备观测、强对抗、实时性与安全约束'],'第2节 问题建模'),
('典型任务链条',['探测/识别：融合雷达、视觉、声学、射频特征','决策/拦截：选择拦截器、航迹与交会策略','协同/博弈：多防御单元对多目标动态分配','电子干扰：功率、频点、波束与时序自适应控制'],'第3节 任务划分'),
('方法谱系总览',['值函数：DQN/Double DQN 适合离散动作决策','策略梯度：PPO/DDPG/SAC 支持连续控制','多智能体RL：处理集群攻防与协同分配','模仿学习：利用专家轨迹降低探索成本','Sim-to-Real：仿真训练后迁移到真实系统'],'第4节 RL方法分类'),
('值函数方法',['面向离散拦截动作、目标选择与资源分配','优势：训练稳定、实现简单、易于解释Q值','局限：连续控制能力弱，状态维度升高后样本需求大','适用：单平台单目标或小规模多目标初步决策'],'方法分类'),
('策略梯度与Actor-Critic',['适合连续机动控制、干扰功率调节与航迹优化','PPO强调稳定更新，SAC强调探索与鲁棒控制','可引入约束奖励处理能耗、碰撞与安全边界','关键在奖励设计和仿真覆盖度'],'方法分类'),
('多智能体强化学习',['适合多拦截器、多传感器与集群反制任务','集中训练、分散执行可兼顾全局协同与部署约束','需处理通信延迟、信用分配和非平稳博弈','潜力方向：对抗自博弈与层级任务分解'],'方法分类'),
('模仿学习与Sim-to-Real',['专家规则/飞手轨迹可用于行为克隆或离线RL','域随机化提升对风场、传感噪声和动力学误差的泛化','仿真平台需覆盖传感器、通信、机动与干扰链路','真实部署需安全屏障与在线监控机制'],'关键技术'),
('代表性工作对比维度',['任务：探测识别、拦截决策、协同攻防、电子干扰','算法：DQN、PPO、DDPG/SAC、MADDPG/QMIX、IL','场景：单目标、多目标、蜂群、复杂地形与遮挡','评价：成功率、响应时间、能耗、鲁棒性、泛化性'],'代表工作'),
('仿真平台与数据集',['AirSim/Unreal：视觉与飞行动力学仿真便利','自研平台：便于定制雷达、电子战和拦截器模型','公开数据不足，常依赖合成轨迹与半实物仿真','数据闭环：仿真—实测—再训练是关键工程路径'],'平台与数据'),
('评测指标体系',['任务有效性：拦截成功率、漏警率、误警率','时效性：发现到决策延迟、交会时间、重规划频率','资源代价：能耗、弹药/干扰资源、通信负载','可信性：鲁棒性、可解释性、安全约束违规率'],'评测指标'),
('系统落地架构',['多源感知层输出目标状态与置信度','决策层融合规则约束与RL策略建议','执行层负责航迹、火控、干扰和协同通信','安全层提供人工接管、策略限幅和日志审计'],'工程落地'),
('局限与挑战',['样本效率低，真实对抗数据难获取','仿真到现实存在动力学、传感器和通信差距','对抗鲁棒性与极端场景覆盖不足','军事/公共安全场景要求高可信与可审计'],'局限挑战'),
('未来研究方向',['安全约束RL与可验证策略学习','多智能体自博弈与博弈论结合','离线RL利用历史演训与仿真数据','跨平台迁移、在线自适应与人机协同决策'],'开放问题'),
('结论',['RL为反无人机系统提供自适应序贯决策能力','短期最可行方向是仿真训练+规则安全约束混合架构','长期突破依赖高保真仿真、真实数据闭环与可信验证','反无人机RL应与传感融合、电子战和指控系统协同设计'],'结论'),
('参考文献与资料来源',['本PPT基于同名中文文献综述DOCX整理','参考原始221条RL anti-UAV主题文献元数据','重点保留综述中的方法分类、任务划分与挑战总结','详细条目请见文献综述末尾参考文献章节'],'参考文献')]
for title,bullets,section in slides: bullet_slide(title,bullets,section)
prs.save(PPT)
print(PPT, PPT.stat().st_size, len(prs.slides))
